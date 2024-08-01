from email.mime.image import MIMEImage
import os
import shutil
import time
import win32com.client
import pyodbc
from pptx import Presentation
import smtplib
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from dotenv import load_dotenv

load_dotenv()

conn_str = os.getenv("STRING_CONNECTION")

def enviar_jpg_por_email(email, caminho_jpg):
    
    mensagem = MIMEMultipart()
    mensagem["From"] = os.getenv("SMTP_EMAIL_REMETENTE")
    mensagem["To"] = email
    mensagem["Subject"] = "Certificado"

    body = "Bom dia!\n\nSegue em anexo seu certificado de participação.\n\nAtenciosamente,\nInovação e Melhoria Contínua"
    mensagem_texto = MIMEText(body)
    mensagem.attach(mensagem_texto)

    # Anexar o JPG à mensagem de e-mail
    with open(caminho_jpg, "rb") as arquivo_jpg:
        anexo_jpg = MIMEImage(arquivo_jpg.read(), _subtype="jpg")
        anexo_jpg.add_header("Content-Disposition", "attachment", filename="Assinatura.jpg")
        mensagem.attach(anexo_jpg)

    # Enviar o e-mail
    with smtplib.SMTP("smtp.gmail.com", 587) as servidor:
        servidor.starttls()
        email = os.getenv("SMTP_EMAIL_REMETENTE")
        senha = os.getenv("SMTP_SENHA_REMETENTE")
        servidor.login(email,senha )
        servidor.send_message(mensagem)


def transformar_em_jpg(caminho_arquivo):
    ppttoJPG = 17
    try:
        powerpoint = win32com.client.Dispatch("Powerpoint.Application")
        deck = powerpoint.Presentations.Open(caminho_arquivo)
        time.sleep(2)
        if caminho_arquivo.endswith(".pptx"):
            deck.SaveAs(caminho_arquivo[:-5], ppttoJPG)
        elif caminho_arquivo.endswith(".ppt"):
            deck.SaveAs(caminho_arquivo[:-4], ppttoJPG)
        deck.Close()
        powerpoint.Quit()   
        print('Salvo em JPG')
        os.remove(caminho_arquivo)
    except:
        print('Não foi possível abrir o arquivo')

def excluir_arquivo(caminho_arquivo):
    if os.path.exists(caminho_arquivo):
        os.remove(caminho_arquivo)
        print("Arquivo excluído com sucesso.")
    else:
        print("O arquivo não existe.")  

def main(id):
    conn = pyodbc.connect(conn_str)
    cursor = conn.cursor()

    query = f"SELECT * FROM Funcionario WHERE idfuncionario = {id}"
    cursor.execute(query)

    funcionario = cursor.fetchone()

    if funcionario:
        nomeArquivo = funcionario.Nome + ".pptx"
        shutil.copy2("Certificado.pptx", nomeArquivo)

        presentation = Presentation(nomeArquivo)

        novo_slide = presentation.slides[0]
        
        for shape in novo_slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        if "NOME" in run.text:
                            run.text = run.text.replace("NOME",  funcionario.Nome)

        presentation.save(nomeArquivo)

        transformar_em_jpg(nomeArquivo)

        caminho_arquivo_jpg = nomeArquivo.replace(".pptx", "") + "\Slide1.JPG"

        enviar_jpg_por_email(funcionario.Email, caminho_arquivo_jpg)
        excluir_arquivo(nomeArquivo.replace(".pptx", ""))
        print("Assinatura enviada com sucesso!")
    else:
        print(f"Informações do funcionário com ID {id} não encontradas.")

main()